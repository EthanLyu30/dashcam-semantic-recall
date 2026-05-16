"""Convert docs/phase2-report/phase2-report.html to phase2-report.pptx.

Strategy:
  1. Parse each <section class="slide"> from the HTML.
  2. Wrap each slide in a standalone HTML page (same CSS, body is exactly
     1280x720 with no scrollbar/padding).
  3. Use Chrome headless to screenshot each page at 1280x720.
  4. Assemble screenshots into a 16:9 python-pptx presentation.
"""
from __future__ import annotations

import re
import subprocess
import sys
import time
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parents[1]
REPORT_DIR = REPO_ROOT / "docs" / "phase2-report"
HTML_FILE  = REPORT_DIR / "phase2-report.html"
OUT_PPTX   = REPORT_DIR / "phase2-report.pptx"
TMP_DIR    = REPORT_DIR / "_slide_tmp"
CHROME     = r"C:\Program Files\Google\Chrome\Application\chrome.exe"

# Widescreen 16:9
SLIDE_W_PX = 1280
SLIDE_H_PX = 720


def extract_css(html: str) -> str:
    m = re.search(r"<style>(.*?)</style>", html, re.DOTALL)
    return m.group(1) if m else ""


def extract_slides(html: str) -> list[str]:
    return re.findall(
        r'<section class="slide"[^>]*>.*?</section>',
        html,
        re.DOTALL,
    )


def build_slide_html(css: str, slide_body: str) -> str:
    return f"""<!doctype html>
<html lang="zh-CN">
<head>
<meta charset="utf-8"/>
<style>
{css}
/* Override: make the slide fill the entire 1280x720 viewport */
html,body{{margin:0;padding:0;width:{SLIDE_W_PX}px;height:{SLIDE_H_PX}px;overflow:hidden;background:var(--bg);}}
.slide{{
  position:relative;width:{SLIDE_W_PX}px;height:{SLIDE_H_PX}px;
  margin:0!important;border-radius:0!important;box-shadow:none!important;
  background:var(--bg);padding:48px 56px 56px 56px;
  display:flex;flex-direction:column;overflow:hidden;
}}
</style>
</head>
<body>
{slide_body}
</body>
</html>"""


def screenshot_slide(html_path: Path, out_png: Path) -> None:
    file_url = html_path.as_uri()
    cmd = [
        CHROME,
        "--headless=new",
        "--disable-gpu",
        "--no-sandbox",
        "--disable-software-rasterizer",
        "--force-device-scale-factor=1",
        "--hide-scrollbars",
        "--run-all-compositor-stages-before-draw",
        f"--window-size={SLIDE_W_PX},{SLIDE_H_PX}",
        f"--screenshot={out_png}",
        file_url,
    ]
    result = subprocess.run(cmd, capture_output=True, timeout=30)
    if result.returncode != 0:
        print(f"  WARN chrome exited {result.returncode}: {result.stderr[:200]}")


def build_pptx(images: list[Path], out: Path) -> None:
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    # 16:9 widescreen in EMU (1 inch = 914400 EMU; 13.33" x 7.5")
    prs.slide_width  = Emu(12192000)   # 13.333... inches
    prs.slide_height = Emu(6858000)    # 7.5 inches

    blank = prs.slide_layouts[6]  # truly blank layout

    for img in images:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(img),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        print(f"  added {img.name}")

    prs.save(str(out))


def main() -> int:
    TMP_DIR.mkdir(parents=True, exist_ok=True)

    html = HTML_FILE.read_text(encoding="utf-8")
    css  = extract_css(html)
    slides = extract_slides(html)
    print(f"Found {len(slides)} slides")

    png_paths: list[Path] = []

    visible_slides = [s for s in slides if 'display:none' not in s[:100]]
    print(f"  ({len(slides) - len(visible_slides)} hidden slide(s) skipped)")
    slides = visible_slides

    for i, slide_body in enumerate(slides, start=1):
        tmp_html = TMP_DIR / f"slide_{i:02d}.html"
        tmp_png  = TMP_DIR / f"slide_{i:02d}.png"

        # Write HTML in the REPORT_DIR so relative image paths resolve
        # (shots/qt-*.png, ../prototype-source/..., etc.)
        real_html = REPORT_DIR / f"_slide_{i:02d}_tmp.html"
        real_html.write_text(build_slide_html(css, slide_body), encoding="utf-8")

        print(f"  screenshotting slide {i:02d} …", end=" ", flush=True)
        screenshot_slide(real_html, tmp_png)

        # Clean up temp HTML
        real_html.unlink(missing_ok=True)

        if tmp_png.exists():
            print(f"ok ({tmp_png.stat().st_size // 1024} KB)")
            png_paths.append(tmp_png)
        else:
            print("FAILED (no png produced)")

    if not png_paths:
        print("ERROR: no screenshots produced — aborting")
        return 1

    print(f"\nAssembling {len(png_paths)} slides into PPTX …")
    build_pptx(png_paths, OUT_PPTX)
    print(f"\nDone → {OUT_PPTX.relative_to(REPO_ROOT)}")

    # Clean up tmp dir
    for f in TMP_DIR.iterdir():
        f.unlink(missing_ok=True)
    TMP_DIR.rmdir()

    return 0


if __name__ == "__main__":
    sys.exit(main())
